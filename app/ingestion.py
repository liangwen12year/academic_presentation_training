"""PPTX ingestion: extract slides as images, speaker notes, and generate scripts."""

from __future__ import annotations

import io
import subprocess
import tempfile
import uuid
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

from app.config import settings


@dataclass
class SlideData:
    index: int
    image_path: str  # relative path served by static files
    notes: str  # speaker notes extracted from PPTX
    script: str  # final script (notes or LLM-generated)
    reference_audio_path: str = ""  # path to golden reference audio
    id: str = field(default_factory=lambda: uuid.uuid4().hex[:8])


@dataclass
class PresentationData:
    id: str
    filename: str
    slides: list[SlideData]


def _render_slides_with_libreoffice(pptx_path: Path, pres_id: str) -> list[str] | None:
    """Render all slides to PNG using LibreOffice headless. Returns list of served paths or None on failure."""
    out_dir = settings.slides_dir / pres_id
    out_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "png",
                    "--outdir",
                    tmp_dir,
                    str(pptx_path),
                ],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
            return None

        # LibreOffice exports one PNG per slide, named like "filename.png" for single-slide
        # or uses a PDF intermediate. For multi-slide, we convert via PDF first.
        png_files = sorted(Path(tmp_dir).glob("*.png"))

        if png_files:
            # Single-slide export: LibreOffice only produces one PNG for the whole file.
            # We need to go through PDF to get per-slide images.
            if len(png_files) == 1:
                # Check if it's truly a single-slide presentation
                prs = Presentation(str(pptx_path))
                if len(prs.slides) > 1:
                    png_files = []  # Fall through to PDF approach

        if not png_files:
            # Use PDF intermediate for per-slide rendering
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        tmp_dir,
                        str(pptx_path),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=120,
                )
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
                return None

            pdf_files = list(Path(tmp_dir).glob("*.pdf"))
            if not pdf_files:
                return None

            pdf_path = pdf_files[0]
            try:
                from pdf2image import convert_from_path

                images = convert_from_path(str(pdf_path), dpi=200)
                paths = []
                for idx, img in enumerate(images):
                    filename = f"slide_{idx:03d}.png"
                    img.save(str(out_dir / filename), "PNG")
                    paths.append(f"/static/slides/{pres_id}/{filename}")
                return paths
            except ImportError:
                # pdf2image not installed, try pdftoppm
                try:
                    subprocess.run(
                        [
                            "pdftoppm",
                            "-png",
                            "-r",
                            "200",
                            str(pdf_path),
                            str(Path(tmp_dir) / "slide"),
                        ],
                        check=True,
                        capture_output=True,
                        timeout=120,
                    )
                    png_files = sorted(Path(tmp_dir).glob("slide-*.png"))
                except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
                    return None

        paths = []
        for idx, png_file in enumerate(png_files):
            filename = f"slide_{idx:03d}.png"
            dest = out_dir / filename
            Image.open(png_file).save(str(dest), "PNG")
            paths.append(f"/static/slides/{pres_id}/{filename}")
        return paths if paths else None


def _extract_slide_image_fallback(prs: Presentation, slide, slide_idx: int, pres_id: str) -> str:
    """Fallback slide image export using python-pptx shape compositing."""
    width = prs.slide_width or Emu(9144000)
    height = prs.slide_height or Emu(6858000)
    px_w = int(width / 9525)
    px_h = int(height / 9525)

    canvas = Image.new("RGB", (px_w, px_h), "#FFFFFF")

    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            img_bytes = shape.image.blob
            try:
                img = Image.open(io.BytesIO(img_bytes))
                left = int(shape.left / 9525) if shape.left else 0
                top = int(shape.top / 9525) if shape.top else 0
                sw = int(shape.width / 9525) if shape.width else img.width
                sh = int(shape.height / 9525) if shape.height else img.height
                img = img.resize((sw, sh), Image.LANCZOS)
                canvas.paste(img, (left, top))
            except Exception:
                pass

    if canvas.getcolors(maxcolors=2) and len(canvas.getcolors(maxcolors=2)) == 1:
        from PIL import ImageDraw, ImageFont

        draw = ImageDraw.Draw(canvas)
        draw.rectangle([(0, 0), (px_w, px_h)], fill="#F0F0F5")
        text = f"Slide {slide_idx + 1}"
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
        except (IOError, OSError):
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        draw.text(((px_w - tw) / 2, (px_h - th) / 2), text, fill="#333355", font=font)

        y_offset = (px_h - th) / 2 + th + 30
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        try:
                            small_font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20
                            )
                        except (IOError, OSError):
                            small_font = ImageFont.load_default()
                        draw.text((40, y_offset), line, fill="#555577", font=small_font)
                        y_offset += 28
                        if y_offset > px_h - 40:
                            break

    out_dir = settings.slides_dir / pres_id
    out_dir.mkdir(parents=True, exist_ok=True)
    filename = f"slide_{slide_idx:03d}.png"
    out_path = out_dir / filename
    canvas.save(str(out_path), "PNG")
    return f"/static/slides/{pres_id}/{filename}"


def _extract_notes(slide) -> str:
    """Extract speaker notes text from a slide."""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        return notes_frame.text.strip()
    return ""


def _extract_slide_text(slide) -> str:
    """Extract all visible text from a slide's shapes."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
    return "\n".join(texts)


async def generate_script_from_slide(slide_text: str, slide_index: int) -> str:
    """Use an LLM to generate a presentation script from slide content."""
    prompt = (
        f"You are a presentation coach. Given the text content of slide {slide_index + 1}, "
        f"write a natural, conversational speaker script (2-4 sentences) that a presenter "
        f"would say while showing this slide. Do not include stage directions or notes — "
        f"just the spoken words. Never use meta-references like 'on this slide', "
        f"'this slide shows', 'as you can see here', or 'let me walk you through this slide'. "
        f"Speak directly about the content as if talking to the audience.\n\n"
        f"Slide content:\n{slide_text}\n\n"
        f"Speaker script:"
    )

    fallback = slide_text if slide_text else f"This is slide {slide_index + 1}."

    def _has_real_key(key: str) -> bool:
        return bool(key) and not key.startswith("your-")

    try:
        if settings.llm_provider == "gemini" and _has_real_key(settings.google_api_key):
            import google.generativeai as genai

            genai.configure(api_key=settings.google_api_key)
            model = genai.GenerativeModel("gemini-1.5-flash")
            response = model.generate_content(prompt)
            return response.text.strip()
        elif _has_real_key(settings.openai_api_key):
            from openai import AsyncOpenAI

            client = AsyncOpenAI(api_key=settings.openai_api_key)
            response = await client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
            )
            return response.choices[0].message.content.strip()
        else:
            return fallback
    except Exception as e:
        print(f"[WARN] LLM script generation failed for slide {slide_index + 1}: {e}")
        return fallback


async def ingest_pptx(filepath: Path) -> PresentationData:
    """Parse a PPTX file and return structured slide data."""
    pres_id = uuid.uuid4().hex[:12]
    prs = Presentation(str(filepath))

    # Try LibreOffice rendering for full-fidelity slide images
    lo_paths = _render_slides_with_libreoffice(filepath, pres_id)
    if lo_paths:
        print(f"[INFO] Rendered {len(lo_paths)} slides with LibreOffice")
    else:
        print("[WARN] LibreOffice rendering failed, using fallback compositing")

    slides: list[SlideData] = []

    for idx, slide in enumerate(prs.slides):
        if lo_paths and idx < len(lo_paths):
            image_path = lo_paths[idx]
        else:
            image_path = _extract_slide_image_fallback(prs, slide, idx, pres_id)

        notes = _extract_notes(slide)

        if notes:
            script = notes
        else:
            slide_text = _extract_slide_text(slide)
            script = await generate_script_from_slide(slide_text, idx)

        slides.append(
            SlideData(
                index=idx,
                image_path=image_path,
                notes=notes,
                script=script,
            )
        )

    return PresentationData(id=pres_id, filename=filepath.name, slides=slides)
